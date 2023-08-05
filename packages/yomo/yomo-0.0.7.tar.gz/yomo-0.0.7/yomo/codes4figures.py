# coding: utf-8
__all__  = ['stamp', 'Look', 'present', 'preserve', 'color']
import re
from datetime import datetime, timezone
from copy import deepcopy
from shutil import copy2
import os
import sys
import inspect
import numpy as np
import numpy.ma as ma
import matplotlib as mpl
import matplotlib.pyplot as plt
import matplotlib.dates as mdates
import matplotlib.ticker as plticker
from matplotlib import cm
# from scipy.stats import gaussian_kde
from netCDF4 import Dataset
from pptx import Presentation # pip install python-pptx 
from pptx.util import Inches  # https://python-pptx.readthedocs.io/
from yomo import index, assemble
def stamp(savefigdir, figname, overwrite='preserve', **keywords):
    """Marks a stamp on the current figure and saves it.
    
    Parameters
    ----------
    savefigdir : string
        A path to the directory under which the figure is to be saved. Do not include the file name.
        Leave it empty ('') and no figure is saved.
    figname : string
        The name of the figure. 
    overwrite : string
        Either 'overwrite', 'backup' or 'no', instructing the action in case the file exists. 
        'backup', the default, takes a copy of the existing file under a slightly modified name and
        saves the new file.
        
    **keywords : optional; see savefig
        Passed on to savefig. 
        
    Returns
    -------
    fname
        A string containing a path to the saved figure. Empty if no figure is saved.

    !!! to be updated
    Raises
    ------
    OVERWRITE FIGURE?
    
    BadException
        Because you shouldn't have done that.
    
    Notes
    -----
    Notes about the implementation algorithm (if needed).
    This can have multiple paragraphs.
    You may include some math:
    .. math:: X(e^{j\omega } ) = x(n)e^{ - j\omega n}
    And even use a greek symbol like :math:`omega` inline.

    Examples
    --------
    These are written in doctest format, and should illustrate how to
    use the function.
    >>> a=[1,2,3]
    >>> print [x + 3 for x in a]
    [4, 5, 6]
    >>> print "a\n\nb"
    a
    b
    !!! end of to be updated
    """
        
    # verify figure name
    if figname:
        stampstr = figname
    else:
        raise ValueError('figname is empty.')        
    if ',' in figname:
        figname = figname[:figname.find(',')]
        
    # hide existing stamps, if any
    fig = plt.gcf()
    for p in set(fig.findobj(lambda x: x.get_gid()=='stamp')):
        p.set_visible(False)

    # determine the caller
    caller = inspect.getframeinfo(sys._getframe(2))[2] # caller = os.path.basename(inspect.stack()[2][1])
    if caller[0] != '<' and caller != 'interactiveshell.py':
        stampstr += ', ' + caller 
        
    # determine the author
    author = os.getenv('USERNAME') # this function might be OS-dependent; if so, update.
    if author:
        stampstr += ', ' + author
    
    # record the time    
    stampstr += ', ' + datetime.now().strftime("%Y-%m-%d %H:%M:%S")

    # stamp now
    plt.figtext(0.01, 0.01, stampstr, size='x-small', gid='stamp')

    # save figure, if requested
    if savefigdir:
        fname = os.path.join(savefigdir, figname)
        preserve(fname, overwrite)
        plt.savefig(fname, bbox_inches='tight', **keywords)
        # plt.savefig(fname, **keywords)
    else:
        fname = ''
        
    # return the path to the figure (if saved) or an empty string 
    return fname
    
class Look:
    """A x-y plot.
    
    Look(dataset, x, ys, savefigdir='', prefix='', overwrite='preserve', show=True, **subsets)
    
    dataset is a NetCDF4 dataset. It must have variables from which x and ys are chosen.
    
    x is a string specifying a single variable that is a matrix or vector.
    
    ys is a string or a list of strings that specify a matrix, a vector or a combination of them.
    
    savefigdir specifies the directory in which each figure is saved. Leave it blank, and no figure is saved. 
    
    prefix precedes an automatically determined string for each figure to form the file name. 
    
    For overwrite, see stamp.
    
    show determines whether the plot is shown. Set this to False if another subplot is to be added to the current figure, for example.
    
    For **subsets, see index.
    """
        
    def __init__(self, dataset, x, ys, xgrid=(), ygrid=(), s='', slim=(), c='', clim=(), cmap=cm.plasma_r, savefigdir='', prefix='', overwrite='preserve', show=True, **subsets):
        # accept not just a dataset but a file as dataset.
        if not hasattr(dataset, 'data_model') or dataset.data_model not in ('NETCDF4', 'NETCDF3_CLASSIC'): 
            try:
                dataset = Dataset(dataset)
            except:
                dataset = assemble(dataset)
        # avoid confusion by a conflict between variable names
        for name in ['xgrid', 'ygrid', 's', 'slim', 'c', 'gaussian_kde', 'clim', 'cmap', 'savefigdir', 'prefix', 'overwrite']:
            if name in dataset.variables.keys():
                raise ValueError('Look cannot handle this dataset because ' + name + ' is among the variables.')
            if name in dataset.dimensions.keys():
                raise ValueError('Look cannot handle this dataset because ' + name + ' is among the dimensions.')
        # open new axes
        self.ax = plt.gca()
        # format y
        if isinstance(ys, str):
            ys = [ys]
        if len(ys)>1 and ma.any(xgrid):
            raise ValueError('ys must not have a single element when xgrid is given.')
        # restrict xgrid and ygrid
        if ma.any(xgrid) and not ma.any(ygrid):
            raise ValueError('xgrid can be given only when ygrid is.')
        if not ma.any(xgrid) and ma.any(ygrid):
            raise ValueError('ygrid can be given only when xgrid is.')
        if ma.any(xgrid) and s:
            raise ValueError('s must not be given when xgrid is.')
        if ma.any(xgrid) and c:
            raise ValueError('c must not be given when xgrid is.')
        # restrict slim
        if len(slim) == 1:
            raise ValueError('slim must be empty or have two or more elements.')
        if slim and not s:
            raise ValueError('slim can be given only when s is.')
        # restrict clim
        if len(clim) == 1:
            raise ValueError('clim must be empty or have two or more elements.')
        if clim and not (c or ma.any(xgrid)):
            raise ValueError('clim can be given only when c or xgrid is.')
        # generate reversed masks and slices
        # Note that the slices almost always contain all information that is contained in reversed masks, but not always. See index.
        variables = [x] + deepcopy(ys)
        if s:
            variables += [s]
        if c:
            variables += [c]
        rmasks, slices = index(dataset, variables=variables, **subsets)
        # get x
        if x=='time_bounds':
            # prepare to plot line segments
            xval = np.array([[datetime.fromtimestamp(t, tz=timezone.utc) for t in ts] for ts in dataset['time_bounds'][:]])
        elif x=='time':
            # Developer's notes:
            # Convert to the format plt.plot accepts. plt.plot_date is not used, as it is slow. 
            # Use datetime.datetime. NetCDF4/num2date is an order of magnitude slower.
            # The one-day difference between Python and matplotlib date notations is properly take care of. 
            # The conversion could be done outside the Look class to enable reuse, but the code would be more complicated and the time saving, negligible. 
            # The notes are valid as of January 2017. Verify as matplotlib and NetCDF4 modules are updated.
            xval = np.array([datetime.fromtimestamp(t, tz=timezone.utc) for t in dataset['time'][:]])
        elif len(dataset[x].dimensions)>=2:
            # See above about the potential need for masking some x elements when x is at least 2-D.
            xval = ma.array(dataset[x][:], mask=(1-rmasks[x]))
        else:
            xval = dataset[x][:]
        xval = xval[np.ix_(*(slices[d] for d in dataset[x].dimensions))]
#         xval = ma.array(xval, mask=1-rmasks[x])[np.ix_(*(slices[d] for d in dataset[x].dimensions))]
        if x=='time_bounds':
            # modify the x from (2,p) to (p*3,) 
            xval = ma.hstack((xval, xval[:,0][:,np.newaxis])).flatten()  
        if xval.ndim>2:
            # plt.plot accepts arrays no greater than 2-D. This is probably the case for the other plotting modules like scatter.
            # Try squeezing.
            xval = np.squeeze(xval)
        # get s
        sc = {'cmap':cmap, 'edgecolors':'none'}
        if s:
            sval = ma.array(dataset[s][:], mask=1-rmasks[s])[np.ix_(*(slices[d] for d in dataset[s].dimensions))]
            # adjust sval so that all markers are within the pre-set size range
            if slim:
                sval[sval<slim[0]]=slim[0]
                sval[sval>slim[-1]]=slim[-1]
            else:
                slim = (ma.min(sval), ma.max(sval))
            slim0 = (3, 20) # the pre-set marker size range. Note that pyplot.scatter takes the square of size as input
            slim0 = (3, 10) # the pre-set marker size range. Note that pyplot.scatter takes the square of size as input
            sc['s'] = ((sval-slim[0])*(slim0[-1]-slim0[0])/(slim[-1]-slim[0]) + slim0[0]) ** 2
        # get c
        if c=='gaussian_kde':
            raise NotImplementedError('gaussian KDE to be implemented.')
            xy = ma.vstack([x,y])
            sc['c'] = gaussian_kde(xy)            
        elif c:
            sc['c'] = ma.array(dataset[c][:], mask=1-rmasks[c])[np.ix_(*(slices[d] for d in dataset[c].dimensions))]
        # get each y and plot it vs x
        for y in ys:
            yval = ma.array(dataset[y][:], mask=1-rmasks[y])[np.ix_(*(slices[d] for d in dataset[y].dimensions))]
            if yval.ndim>2:
                # plt.plot accepts arrays no greater than 2-D. This is probably the case for the other plotting modules like scatter.
                # Try squeezing.
                yval = np.squeeze(yval)
            if dataset[y].dimensions[0]!=dataset[x].dimensions[0] and len(dataset[y].dimensions)>1 and dataset[y].dimensions[1]==dataset[x].dimensions[0]:
                yval = yval.T
            if x=='time_bounds': # yes, it's x.
                # prepare to plot line segments
                yval = yval.repeat(3,axis=0)            
                yval.mask[::3] = True
            # plot
            if ma.any(xgrid) and len(clim)==2: # color density plot
                _,_,_,linesi = self.ax.hist2d(xval, yval, bins=[xgrid, ygrid], cmin=clim[0], cmax=clim[1], cmap=cmap)
            elif ma.any(xgrid) and not clim: # color density plot
                _,_,_,linesi = self.ax.hist2d(xval, yval, bins=[xgrid, ygrid], cmap=cmap)
            elif ma.any(xgrid) and len(clim)>2:
                raise NotImplementedError('Layered color plots to be implemented.')
            elif (s or c) and len(clim)<=2: # scatter plot
                # Note the _r in the color map, making the map go from light colors to dense ones. 
                # This is intuitive for quantities like AOD and CO_ppbv. 
                # For temperature, switch map manually. 
                self.ax.autoscale(True, tight=True)
                linesi = self.ax.scatter(xval, yval, **sc)
                if len(clim)==2:
                    linesi.set_clim(clim)
            elif (s or c) and (len(slim)>2 or len(clim)>2):
                # !!! allow scatter to be created with incremental size bounds or color bounds - faster to make and easier to view for some cases.
                raise NotImplementedError('Layered color plots to be implemented.')
            elif x=='time_bounds':
                linesi = self.ax.plot(xval, yval, '.', linestyle='-', linewidth=2, label=dataset[y].name.replace('_',' '))
            else: # x-y plot
                linesi = self.ax.plot(xval, yval, '.', label=dataset[y].name.replace('_',' '))
            if 'lines' in locals():
                # !!! This will cause an error if linesi is not iterable, as is the case for scatter plots. 
                lines.extend(linesi)
            else:
                lines = linesi
                
        # label
        plt.title(prefix.replace('_',' '))        
        if x=='time' or x=='time_bounds':
            # !!! to-do: look for tzinfo and verify it's UTC.
            xlim = plt.gca().get_xlim()
            if xlim[-1] - xlim[0] <= 2:
                xlabel = 'UTC'                
                # plt.gca().xaxis.set_major_formatter(mdates.DateFormatter('%H'))
                # loc = plticker.MultipleLocator(base=max(np.round((xlim[-1] - xlim[0])*3)/24., 1/24.)) # used to be: plt.gca().xaxis.set_major_locator(mdates.HourLocator()) 
                # plt.gca().xaxis.set_major_locator(loc)
                plt.gcf().autofmt_xdate() # See if we like this. If not, revert to the lines above.
            else:
                mdmd = [mdates.num2date(xlim[0]).year, mdates.num2date(xlim[0]).strftime('%B'), 
                            mdates.num2date(xlim[1]).year, mdates.num2date(xlim[1]).strftime('%B')]
                xlabel = str(mdmd[0])+' '+mdmd[1]
                if mdmd[0] != mdmd[2]:
                    xlabel = xlabel+' - '+str(mdmd[2])+' '+mdmd[3]
                elif mdmd[1] != mdmd[3]:
                    xlabel = xlabel+' - '+mdmd[3]
                plt.gca().xaxis.set_major_formatter(mdates.DateFormatter('%d'))
        else:
            xlabel = dataset[x].long_name if 'long_name' in dataset[x].ncattrs() and re.sub(r'\s+','', dataset[x].long_name) else dataset[x].name
            xlabel = xlabel + ', ' + dataset[x].units if 'units' in dataset[x].ncattrs() and dataset[x].units else xlabel                
        plt.xlabel(xlabel.replace('_',' '))
        ylabel = '' # Note that ys can be plural unlike x.
        for y in ys:
            # legend labels for a special case
            if dataset[y].dimensions[0]!=dataset[x].dimensions[0] and len(dataset[y].dimensions)>1 and dataset[y].dimensions[1]==dataset[x].dimensions[0] and len(dataset[y].dimensions)>=2 and dataset[y].dimensions[0]=='time' and not c:
                [lines[i].set_label(datetime.fromtimestamp(t, tz=timezone.utc)) for i, t in enumerate(dataset['time'][:][slices['time']])]
            # y axis label
            if 'long_name' in dataset[y].ncattrs() and re.sub(r'\s+','', dataset[y].long_name) :
                ylabel = ylabel + dataset[y].long_name + ', '
            else:
                ylabel = ylabel + dataset[y].name + ', '
            if 'units' in dataset[y].ncattrs() and dataset[y].units:
                ylabel = ylabel + dataset[y].units + ', '
        plt.ylabel(ylabel.replace('_',' ')[:-2])
        if s:
            xylims = [self.ax.get_xlim(), self.ax.get_ylim()]
            slabel = dataset[s].long_name if 'long_name' in dataset[s].ncattrs() and re.sub(r'\s+','', dataset[s].long_name) else dataset[s].name
            slabel = slabel + ', ' + dataset[s].units if 'units' in dataset[s].ncattrs() and dataset[s].units else slabel
            sc0 = deepcopy(sc)
            if c:
                sc0['c'] = ma.mean(sc['c'])
            sc1 = deepcopy(sc0)
            sc2 = deepcopy(sc0)
            sc0['s'] = slim0[0]**2
            sc1['s'] = ma.mean(slim0)**2
            sc2['s'] = slim0[-1]**2
            try:
                markersforlegend0 = plt.scatter([],[], **sc0)
                markersforlegend1 = plt.scatter([],[], **sc1)
                markersforlegend2 = plt.scatter([],[], **sc2)
            except:
                try:
                    markersforlegend0 = plt.scatter(None,None, **sc0)
                    markersforlegend1 = plt.scatter(None,None, **sc1)
                    markersforlegend2 = plt.scatter(None,None, **sc2)
                except:
                    markersforlegend0 = plt.scatter(np.array([]).astype(type(xval[0])),np.array([]).astype(type(yval[0])), **sc0)
                    markersforlegend1 = plt.scatter(np.array([]).astype(type(xval[0])),np.array([]).astype(type(yval[0])), **sc1)
                    markersforlegend2 = plt.scatter(np.array([]).astype(type(xval[0])),np.array([]).astype(type(yval[0])), **sc2)
            # !!! there must be an empty variable that is compatible with all sorts of plots (scatter, time on x axis, etc.)
            self.ax.set_xlim(xylims[0])
            self.ax.set_ylim(xylims[1])
            legend = self.ax.legend((markersforlegend0,markersforlegend1,markersforlegend2), 
                       [str(slim[0]), str(ma.mean(slim).astype(type(slim[0]))), str(slim[-1])], title=slabel.replace('_',' '), ncol=3)
            plt.setp(legend.get_title(),fontsize=plt.rcParams['legend.fontsize'])
        if len(dataset[x].dimensions)>=2 and 'wavelength' in dataset[x].dimensions[1] and dataset[x].dimensions[1] in dataset.variables.keys(): 
            color(lines, dataset[dataset[x].dimensions[1]][:][slices[dataset[x].dimensions[1]]], colorbar=True)
        elif len(dataset[y].dimensions)>=2 and 'wavelength' in dataset[y].dimensions[1] and dataset[y].dimensions[1] in dataset.variables.keys() and len(dataset[dataset[y].dimensions[1]][:][slices[dataset[y].dimensions[1]]])==len(lines):
            color(lines, dataset[dataset[y].dimensions[1]][:][slices[dataset[y].dimensions[1]]], colorbar=True)
        elif c: 
            cbar = plt.colorbar(lines)
            if c=='time':
                # !!! to-do: look for tzinfo and verify it's UTC.
                clim = lines.get_clim()
                if clim[1] - clim[0] <= 2:
                    clabel = 'UTC'
                    cbar.ax.yaxis.set_major_locator(mdates.HourLocator())
                    cbar.ax.yaxis.set_major_formatter(mdates.DateFormatter('%H:%M'))
                else:
                    mdmd = [datetime.fromtimestamp(clim[0], tz=timezone.utc).year, datetime.fromtimestamp(clim[0], tz=timezone.utc).strftime('%B'), 
                            datetime.fromtimestamp(clim[1], tz=timezone.utc).year, datetime.fromtimestamp(clim[1], tz=timezone.utc).strftime('%B')]
                    clabel = str(mdmd[0])+' '+mdmd[1]
                    if mdmd[0] != mdmd[2]:
                        clabel = clabel+' - '+str(mdmd[2])+' '+mdmd[3]
                    elif mdmd[1] != mdmd[3]:
                        clabel = clabel+' - '+mdmd[3]
                    cbar.ax.set_ylabel(clabel)
    #                     cbar.ax.yaxis.set_major_formatter(mdates.DateFormatter('%d')) # !!! this causes an error about ordinal.
                    yt = np.arange((clim[1]-clim[0])/86400)*86400 + clim[0]             
                    cbar.set_ticks(yt)
                    cbar.set_ticklabels([datetime.fromtimestamp(t, tz=timezone.utc).day for t in yt])
            else:
                clabel = dataset[c].long_name if 'long_name' in dataset[c].ncattrs() and re.sub(r'\s+','', dataset[c].long_name) else dataset[c].name
                clabel = clabel + ', ' + dataset[c].units if 'units' in dataset[c].ncattrs() and dataset[c].units else clabel                
                cbar.ax.set_ylabel(clabel.replace('_',' '))
        elif ma.any(xgrid): 
            cbar = plt.colorbar(lines)
            cbar.ax.set_ylabel('Occurrence')
        elif s:
            pass
        elif (len(ys)>1 or len(lines)>1) and len(lines)<=10:
            legend = self.ax.legend(loc='upper right')
            plt.setp(legend.get_title(),fontsize=plt.rcParams['legend.fontsize'])
        # prepare to save
        self.lines = lines
        self.savefigdir = savefigdir
        self.prefix = prefix
        self.fignamemain = ''.join([dataset[y].name for y in ys]) + 'vs' + dataset[x].name
        self.fignamemain = self.fignamemain + 'occurrence' if ma.any(xgrid) else self.fignamemain
        self.fignamemain = self.fignamemain + 'by' + dataset[s].name if s else self.fignamemain
        self.fignamemain = self.fignamemain + 'with' + dataset[c].name if c else self.fignamemain
        self.overwrite = overwrite    
        self.show = show
            
    def __enter__(self):
        return self
    
    def __exit__(self, exc_type, exc_value, traceback):
        # save and show
        if self.savefigdir:
            figname = self.prefix + '_' + self.fignamemain + '.png' if self.prefix else self.fignamemain + '.png'
            stamp(self.savefigdir, figname, self.overwrite)
        if self.show:
            plt.show()
        
        # !!! I don't fully understand the lines below.
        # !!! copied from http://stackoverflow.com/questions/22417323/how-do-enter-and-exit-work-in-python-decorator-classes
        if exc_type is not None:
            print(exc_type)
            print(exc_value)
            print(traceback)
        # return False # uncomment to pass exception through

        return self
def present(contents, savepptdir='', pptname='caller'):
    """Generates a PowerPoint presentation with images.

    Parameters
    ----------
    contents : list, tuple  
        Each element must contain one string or more. An integer may be included at the end, or the code assigns an integer for each element. 
        Each string must be either the path to an image or a text. 
        Each integer must specify the slide layout: 0 or 5 for title only, 2 for section header, 6.1 for an image on a blank slide, 6.2 for up to four images on a blank slide, 6.4 for up to four images on a blank slide. 
    savepptdir : string
        A path to the directory under which the presentation is to be saved. Do not include the file name.
        Leave it empty ('') and no presentation is saved.
    pptname : string
        The name of the presentation.
        
    Returns
    -------
    fname
        A string containing a path to the saved presentation. Empty if no presentation is saved.

    !!! to be updated
    Raises
    ------
    OVERWRITE FIGURE?
    
    BadException
        Because you shouldn't have done that.
    !!! end of to be updated
    
    Notes
    -----
    This code requires python-pptx http://python-pptx.readthedocs.io/en/latest/. Use it for creating a presentation with non-image contents.

    !!! to be updated
    Examples
    --------
    These are written in doctest format, and should illustrate how to
    use the function.
    >>> a=[1,2,3]
    >>> print [x + 3 for x in a]
    [4, 5, 6]
    >>> print "a\n\nb"
    a
    b
    !!! end of to be updated
    """
    
    # determine the layouts
    if isinstance(contents, (list, tuple)):
        layouts = [i[-1] if isinstance(i, (list,tuple)) and isinstance(i[-1], (int, float)) else 6+len(i)/10 if isinstance(i, (list,tuple)) and not isinstance(i[-1], (int, float)) else 6.1 if os.path.exists(i) else 0 for i in contents]
        contents = [i[:-1] if isinstance(i, (list,tuple)) and isinstance(i[-1], (int, float)) else i for i in contents]
    else:
        raise ValueError('contents must be either a list, tuple, dictionary or OrderedDict.')        
            
    # open a blank presentation
    prs = Presentation()

    # add images or a title
    previouslayout = -1
    for content, layout in zip(contents, layouts):
        if layout // 1.0 == 6.0: # layout #6 comes with multiple options
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            if layout == 6 or layout == 6.1: # one image
                width = Inches(10)
                height = Inches(7.5)
                left = top = Inches(0)
                previouslayout = 6.1
            elif layout == 6.2: # four images in one slide
                top = Inches(1.5)
                width = Inches(5)
                height = Inches(3.75)
                if previouslayout == 6.21:
                    left = Inches(5)
                    previouslayout = 6.22
                else:
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    left = Inches(0)
                    previouslayout = 6.21
            elif layout == 6.3 | layout == 6.4: # four images in one slide
                width = Inches(5)
                height = Inches(3.75)
                if previouslayout == 6.41:
                    left = Inches(5)
                    top = Inches(0)
                    previouslayout = 6.42
                elif previouslayout == 6.42:
                    left = Inches(0)
                    top = Inches(3.75)
                    previouslayout = 6.43
                elif previouslayout == 6.43:
                    left = Inches(5)
                    top = Inches(3.75)
                    previouslayout = 6.44
                else:
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    left = top = Inches(0)
                    previouslayout = 6.41
            pic = slide.shapes.add_picture(content, left, top, width=width, height=height)                
        else: 
            slide = prs.slides.add_slide(prs.slide_layouts[layout])
            if layout == 0 or layout == 2 or layout == 5: # assume the content is a title
                title = slide.shapes.title
                title.text = content
            else:
                raise NotImplementedError('Update imgs2ppt.')
            previouslayout = layout
            
    # save the presentation, if requested
    if savepptdir:
        # determine the file name
        if pptname == 'caller': 
            caller = inspect.getframeinfo(sys._getframe(1))[2] # caller = os.path.basename(inspect.stack()[1][1])
            if caller[0] != '<':
                pptname = os.path.splitext(caller)[0] + '.ppt'
            else:
                raise NameError('Caller is not defined. Specify pptname.')
        fname = os.path.join(savepptdir, pptname)
        preserve(fname, 'preserve')
        prs.save(fname)
    else:
        # !!! show the ppt file, when saving is not requested.
        raise NotImplementedError('Modify imgs2ppt to show the presentation.')
        fname = ''

    # return the path to the presentation (if saved) or an empty string 
    return fname
    
def preserve(savefname, overwrite='preserve', replace={}):
    
    if os.path.exists(savefname):
        # verify the instructions
        if overwrite=='preserve':
            preservedfname = savefname.replace('.', '_version' + datetime.fromtimestamp(os.path.getmtime(savefname)).strftime("%Y%m%dT%H%M%S") + '.')
            if replace:
                os.rename(savefname,preservedfname)
                f = open(preservedfname,'r')
                filedata = f.read()
                f.close()
                for key, val in replace.items():
                    filedata = filedata.replace(key, val)
                f = open(savefname,'w')
                f.write(filedata)
                f.close()
            else:
                copy2(savefname, preservedfname)
        elif not overwrite:
            raise FileExistsError(savefname + ' already exists.')
    
def color(objs, wavelengths, colorbar=True, fmt=None):
    """Colors lines according to the wavelength.
    """
        
    # get the unit and format of the wavelengths
    if min(wavelengths)>=100: # assume nm as unit, not um
        wavelengths=wavelengths/1000
        unit='nm'
        if not fmt:
            fmt='{:0.1f}'
    else:
        unit='\mum'
        if not fmt:
            fmt='{:0.3f}'
    # ensure consistency between input variables
    nw = len(wavelengths)
    if len(objs)!=nw and nw!=1:
        raise ValueError('Inconsistent numbers of elements between plot_axes and wavelengths.')
    # load visible color data
    wavelengths0, cm0 = _loadpresetspectrumcolormap()
    # relate the wavelengths to the pre-set colors
    cm = []
    for i in range(3): 
        cm.append(np.interp(wavelengths, wavelengths0, cm0[:,i]))
    cm = np.array(cm).T
    # update the colors of the objects
    legend_string=[]
    for i, obj in enumerate(objs):
        if nw==1:
            i=0 # a tweak for cases where a single color should be applied to multiple lines
        plt.setp(obj, color=cm[i])
        legend_string.append([fmt.format(wavelengths[i])+' '+unit])        
    # put a colorbar if requested
    if colorbar:
        gca = plt.gca()
        gcf = plt.gcf()
        pos = gca.get_position()
        gca.set_position([pos.x0, pos.y0,  pos.width-0.08, pos.height])
        ax2 = gcf.add_axes([pos.x0+pos.width-0.08+0.01, pos.y0, 0.07, pos.height])
        cmap = mpl.colors.ListedColormap(np.flipud(cm))
        bounds = np.arange(nw+1)
        norm = mpl.colors.BoundaryNorm(bounds, cmap.N)
        mpl.colorbar.ColorbarBase(ax2, cmap=cmap, ticks=[], norm=norm,)
        if nw==1:
            i=0
            ax2.text(0.5, 1-(i+0.5)/nw, str(int(np.round(wavelengths[i]*1000))), color='w', horizontalalignment='center', verticalalignment='center', size=16, weight='bold')
        else:
            for i, obj in enumerate(objs):
                ax2.text(0.5, 1-(i+0.5)/nw, str(int(np.round(wavelengths[i]*1000))), color='w', horizontalalignment='center', verticalalignment='center', size=16, weight='bold')
        plt.axes(gca)
        
def _loadpresetspectrumcolormap():
    """Load the color map determined in setspectrumcolor.m."""

    import numpy as np

    # load the wavelengths
    interval=10
    thr=[0, 171, 360, 490, 550, 730, 900, 1100, 2200]
    nuv=np.arange(thr[1],thr[2],interval)/1000
    vis2=np.arange(thr[2],thr[3],interval)/1000
    vis3=np.arange(thr[3],thr[4],interval)/1000
    vis4=np.arange(thr[4],thr[5],interval)/1000
    nir1=np.arange(thr[5],thr[6],interval)/1000
    nir2=np.arange(thr[6],thr[7],interval)/1000
    ir=np.arange(thr[7],thr[8],interval)/1000
    wavelengths0=np.concatenate([nuv,vis2,vis3,vis4,nir1,nir2,ir])
    # load the color map
    cm0 = np.array([[0.50147, 0, 0.50147],
            [0.53079, 0, 0.53079],
            [0.56012, 0, 0.56012],
            [0.58944, 0, 0.58944],
            [0.61877, 0, 0.61877],
            [0.64809, 0, 0.64809],
            [0.67742, 0, 0.67742],
            [0.70674, 0, 0.70674],
            [0.73607, 0, 0.73607],
            [0.7654, 0, 0.7654],
            [0.79472, 0, 0.79472],
            [0.82405, 0, 0.82405],
            [0.85337, 0, 0.85337],
            [0.8827, 0, 0.8827],
            [0.91202, 0, 0.91202],
            [0.94135, 0, 0.94135],
            [0.97067, 0, 0.97067],
            [1, 0, 1],
            [0, 0, 0.7],
            [0, 0, 0.8],
            [0, 0, 0.9],
            [0, 0, 1],
            [0, 0.066667, 0.95],
            [0, 0.13333, 0.9],
            [0, 0.2, 0.85],
            [0, 0.26667, 0.8],
            [0, 0.33333, 0.75],
            [0, 0.4, 0.7],
            [0, 0.46667, 0.65],
            [0, 0.53333, 0.6],
            [0, 0.6, 0.55],
            [0, 0.68889, 0],
            [0, 0.57778, 0],
            [0, 0.46667, 0],
            [0, 0.35556, 0],
            [0, 0.24444, 0],
            [0, 0.13333, 0],
            [0.1, 0.66667, 0],
            [0.2, 0.66667, 0],
            [0.3, 0.66667, 0],
            [0.4, 0.66667, 0],
            [0.5, 0.66667, 0],
            [0.55, 0.6, 0],
            [0.6, 0.53333, 0],
            [0.65, 0.46667, 0],
            [0.7, 0.4, 0],
            [0.75, 0.33333, 0],
            [0.8, 0.26667, 0],
            [0.85, 0.2, 0],
            [0.9, 0.13333, 0],
            [0.95, 0.066667, 0],
            [1, 0, 0],
            [0.9, 0, 0],
            [0.8, 0, 0],
            [0.7, 0, 0],
            [0, 0, 0],
            [0.053712, 0.03255, 0.0066812],
            [0.10742, 0.0651, 0.013362],
            [0.16114, 0.09765, 0.020044],
            [0.21485, 0.1302, 0.026725],
            [0.26856, 0.16275, 0.033406],
            [0.32227, 0.1953, 0.040087],
            [0.37599, 0.22785, 0.046769],
            [0.4297, 0.2604, 0.05345],
            [0.48341, 0.29295, 0.060131],
            [0.53712, 0.3255, 0.066812],
            [0.59084, 0.35805, 0.073494],
            [0.64455, 0.3906, 0.080175],
            [0.68264, 0.42315, 0.086856],
            [0.65823, 0.4557, 0.093537],
            [0.63381, 0.48825, 0.10022],
            [0.6094, 0.5208, 0.1069],
            [0.1953, 0.40627, 0.44655],
            [0.18502, 0.41997, 0.44936],
            [0.17474, 0.43368, 0.45218],
            [0.16446, 0.44738, 0.45499],
            [0.16076, 0.46109, 0.4578],
            [0.18338, 0.47479, 0.46062],
            [0.20599, 0.4885, 0.46343],
            [0.22861, 0.5022, 0.46624],
            [0.25123, 0.51591, 0.46906],
            [0.27384, 0.52961, 0.47187],
            [0.29646, 0.54332, 0.47468],
            [0.31907, 0.55702, 0.47749],
            [0.34169, 0.57073, 0.48031],
            [0.36431, 0.58444, 0.48312],
            [0.38692, 0.59814, 0.48593],
            [0.40954, 0.61185, 0.48875],
            [0.43215, 0.62555, 0.49156],
            [0.45477, 0.63926, 0.49437],
            [0.47738, 0.65296, 0.49719],
            [0.5, 0.66667, 0.5],
            [0.23214, 0.66667, 0.5],
            [0.21429, 0.66667, 0.5],
            [0.19643, 0.66667, 0.5],
            [0.17857, 0.66667, 0.5],
            [0.16071, 0.66667, 0.5],
            [0.14286, 0.66667, 0.5],
            [0.125, 0.66667, 0.5],
            [0.10714, 0.66667, 0.5],
            [0.089286, 0.66667, 0.5],
            [0.071429, 0.66667, 0.5],
            [0.053571, 0.66667, 0.5],
            [0.035714, 0.66667, 0.5],
            [0.017857, 0.66667, 0.5],
            [0, 0.66667, 0.5],
            [0.0089286, 0.65476, 0.50893],
            [0.017857, 0.64286, 0.51786],
            [0.026786, 0.63095, 0.52679],
            [0.035714, 0.61905, 0.53571],
            [0.044643, 0.60714, 0.54464],
            [0.053571, 0.59524, 0.55357],
            [0.0625, 0.58333, 0.5625],
            [0.071429, 0.57143, 0.57143],
            [0.080357, 0.55952, 0.58036],
            [0.089286, 0.54762, 0.58929],
            [0.098214, 0.53571, 0.59821],
            [0.10714, 0.52381, 0.60714],
            [0.11607, 0.5119, 0.61607],
            [0.125, 0.5, 0.625],
            [0.13393, 0.4881, 0.63393],
            [0.14286, 0.47619, 0.64286],
            [0.15179, 0.46429, 0.65179],
            [0.16071, 0.45238, 0.66071],
            [0.16964, 0.44048, 0.66964],
            [0.17857, 0.42857, 0.67857],
            [0.1875, 0.41667, 0.6875],
            [0.19643, 0.40476, 0.69643],
            [0.20536, 0.39286, 0.70536],
            [0.21429, 0.38095, 0.71429],
            [0.22321, 0.36905, 0.72321],
            [0.23214, 0.35714, 0.73214],
            [0.24107, 0.34524, 0.74107],
            [0.25, 0.33333, 0.75],
            [0.26786, 0.33333, 0.73214],
            [0.28571, 0.33333, 0.71429],
            [0.30357, 0.33333, 0.69643],
            [0.32143, 0.33333, 0.67857],
            [0.33929, 0.33333, 0.66071],
            [0.35714, 0.33333, 0.64286],
            [0.375, 0.33333, 0.625],
            [0.39286, 0.33333, 0.60714],
            [0.41071, 0.33333, 0.58929],
            [0.42857, 0.33333, 0.57143],
            [0.44643, 0.33333, 0.55357],
            [0.46429, 0.33333, 0.53571],
            [0.48214, 0.33333, 0.51786],
            [0.5, 0.33333, 0.5],
            [0.51786, 0.33333, 0.48214],
            [0.53571, 0.33333, 0.46429],
            [0.55357, 0.33333, 0.44643],
            [0.57143, 0.33333, 0.42857],
            [0.58929, 0.33333, 0.41071],
            [0.60714, 0.33333, 0.39286],
            [0.625, 0.33333, 0.375],
            [0.64286, 0.33333, 0.35714],
            [0.66071, 0.33333, 0.33929],
            [0.67857, 0.33333, 0.32143],
            [0.69643, 0.33333, 0.30357],
            [0.71429, 0.33333, 0.28571],
            [0.73214, 0.33333, 0.26786],
            [0.75, 0.33333, 0.25],
            [0.74107, 0.34524, 0.24107],
            [0.73214, 0.35714, 0.23214],
            [0.72321, 0.36905, 0.22321],
            [0.71429, 0.38095, 0.21429],
            [0.70536, 0.39286, 0.20536],
            [0.69643, 0.40476, 0.19643],
            [0.6875, 0.41667, 0.1875],
            [0.67857, 0.42857, 0.17857],
            [0.66964, 0.44048, 0.16964],
            [0.66071, 0.45238, 0.16071],
            [0.65179, 0.46429, 0.15179],
            [0.64286, 0.47619, 0.14286],
            [0.63393, 0.4881, 0.13393],
            [0.625, 0.5, 0.125],
            [0.61607, 0.5119, 0.11607],
            [0.60714, 0.52381, 0.10714],
            [0.59821, 0.53571, 0.098214],
            [0.58929, 0.54762, 0.089286],
            [0.58036, 0.55952, 0.080357],
            [0.57143, 0.57143, 0.071429],
            [0.5625, 0.58333, 0.0625],
            [0.55357, 0.59524, 0.053571],
            [0.54464, 0.60714, 0.044643],
            [0.53571, 0.61905, 0.035714],
            [0.52679, 0.63095, 0.026786],
            [0.51786, 0.64286, 0.017857],
            [0.50893, 0.65476, 0.0089286],
            [0.5, 0.66667, 0],
            [0.5, 0.66667, 0.017857],
            [0.5, 0.66667, 0.035714],
            [0.5, 0.66667, 0.053571],
            [0.5, 0.66667, 0.071429],
            [0.5, 0.66667, 0.089286],
            [0.5, 0.66667, 0.10714],
            [0.5, 0.66667, 0.125],
            [0.5, 0.66667, 0.14286],
            [0.5, 0.66667, 0.16071],
            [0.5, 0.66667, 0.17857],
            [0.5, 0.66667, 0.19643],
            [0.5, 0.66667, 0.21429],
            [0.5, 0.66667, 0.23214]])
    return wavelengths0, cm0
